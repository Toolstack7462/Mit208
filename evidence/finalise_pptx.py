"""Finalise the PhishGuard technical presentation by editing the existing deck.

This does not recreate the deck. All ten slides, their layout, colours, fonts and
structure are kept. It makes five kinds of change:

  1. Replaces every superseded figure with the verified one from the 11 August 2026
     run, and corrects two claims the code does not support — the deck said there
     was no rate limiting (there is a per-IP failed-login limiter) and that the
     public CI result was still to be confirmed (it passes).
  2. Removes the placeholder student fields and the stale "release link" line,
     leaving ruled fill-in fields and the tag that actually exists.
  3. Rebuilds slide 2 so the role badges align with their card titles instead of
     floating in the middle of the card, and fills the empty lower half of each
     card with the concrete needs the system meets.
  4. Rebuilds slide 8. Seven screenshots at 2.85 x 1.78in reduced a 1600px-wide
     interface to about 2pt of text — unreadable from a seat. It now shows four
     crops, laid out by aspect ratio, where the interface text renders at roughly
     1.1% of slide height: legible when projected.
  5. Swaps in the simplified layered architecture drawing and updates the speaker
     notes to match every figure on the slide.

Every text substitution is checked: if an expected string is not found the script
fails rather than silently leaving a stale figure on a slide.

Usage:
    python evidence/finalise_pptx.py <input.pptx> <output.pptx>
"""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
SHOTS = HERE / "screenshots"
FIGURES = HERE / "figures"
TMP = HERE / "_pptx_tmp"

TAG = "v1.4-final"
REPO = "github.com/Toolstack7462/Mit208"

# The presenter's prompt sheet lands here, outside the submission folder.
REHEARSAL_DIR = Path(r"D:\Saif Assignement\Private Rehearsal")

# Deck palette and type, read off the existing shapes.
NAVY = RGBColor(0x0B, 0x1F, 0x33)
MUTED = RGBColor(0x52, 0x60, 0x6D)
BRAND = RGBColor(0x25, 0x63, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Aptos"

# Short enough that name and ID stay on one line in the title slide's text box.
BLANK = "_" * 14

# ---------------------------------------------------------------------------
# Verified figures, 11 August 2026 (docs/TESTING.md)
# ---------------------------------------------------------------------------
BACKEND_TESTS = "170"
COVERAGE_PCT = "90%"
COVERAGE_FRAC = "847 / 943"
SMOKE = "22/22"
FRONTEND_TESTS = "92/92"
TOTAL_TESTS = "262"

# Crop boxes shared with the report figure, so the deck and the report show the
# same four steps framed identically. Kept in step with
# evidence/make_workflow_figure.py.
PANES = [
    ("04-analyst-inbox.png", (585, 235, 1465, 1075), "1", "Triage: ordered by risk"),
    ("06-email-detail-explainable-score.png", (1530, 690, 2735, 1665), "2",
     "Explain: a named indicator per point"),
    ("10-release-request-validation.png", (1176, 680, 2024, 1328), "3",
     "Challenge: justified release request"),
    ("13-release-request-approved.png", (575, 245, 3150, 885), "4",
     "Decide: analyst approves, email released and audited in one transaction"),
]

# (old, new) applied to every run in the deck. Each must match at least once.
REPLACEMENTS: list[tuple[str, str]] = [
    # --- Slide 1: identity -------------------------------------------------
    ("[Student name]   |   [Student ID]",
     f"Student name {BLANK}   |   Student ID {BLANK}"),
    ("Final release: [insert genuine v1.0-final link]",
     f"Assessed version: tag {TAG}"),
    ("Final repository: Toolstack7462/Mit208", f"Repository: {REPO}"),
    ("LIVE APPLICATION - CAPTURED 5 AUG 2026", "LIVE APPLICATION - CAPTURED 11 AUG 2026"),
    ("Prepared 5 August 2026", "Prepared 12 August 2026"),

    # --- Slide 3: the rate-limiting claim was wrong ------------------------
    # A per-IP failed-login limiter is implemented in app/ratelimit.py. What the
    # prototype lacks is a limiter that holds state across processes.
    ("No production MFA, rate limiting or token revocation",
     "No MFA, token revocation or cross-process rate limiting"),

    # --- Slide 4: the flow line wrapped onto a second line that sat below its
    # container. The four colour chips are handled by shape name in CHIPS, because
    # "API" and "DATA" as bare strings would match a dozen places in the deck.
    ("Bearer JWT -> FastAPI routers -> validation / RBAC / state rules -> "
     "SQLAlchemy transaction -> PostgreSQL target",
     "Bearer JWT -> routers -> validation / RBAC / state rules -> SQLAlchemy -> PostgreSQL"),

    # --- Slide 5 -----------------------------------------------------------
    ("Node policy/helper tests + CI definition",
     "vitest component tests + CI on every push"),

    # --- Slide 7: verified figures ----------------------------------------
    ("Verified on SQLite and PostgreSQL 16.6, with manual gaps kept visible",
     "Verified on SQLite and PostgreSQL 16.6, with the limitations stated"),
    ("Verified on PostgreSQL 16.6: schema, constraints, full suite and live workflow. "
     "Still manual: confirming the public CI result and recording the narration.",
     "Verified on PostgreSQL 16.6: schema applied, constraints confirmed, full suite and "
     "live workflow passed, and the same suite green in continuous integration."),

    # --- Slide 8: describe the evidence, not the tool that produced it ------
    ("Captured from the running application with Playwright (3200 x 2000)",
     "Screenshots from the running application"),

    # --- Slide 9 -----------------------------------------------------------
    ("Limitations / not yet final", "Limitations"),
    ("Public CI result still to confirm; PostgreSQL verified locally",
     "Small synthetic dataset; no measured behaviour at production volume"),
    # "yet" reads as provisional in a submitted deck.
    ("No automated browser-level end-to-end test yet",
     "No automated browser-level end-to-end test"),
    # docs/ now holds eight: API, ARCHITECTURE, BUG_LOG, DEMO, ERD, SECURITY,
    # SUBMISSION_CHECKLIST and TESTING.
    ("Repository holds setup, schema, six technical documents and evidence",
     "Repository holds setup, schema, eight technical documents and evidence"),
    ("Evaluation: strong local technical evidence; submission completeness depends on "
     "final reproducibility and demonstration.",
     "Evaluation: the core MVP is complete, verified on both database engines and "
     "reproducible from the repository, with its limitations stated rather than implied."),

    # Slide 10 is rebuilt wholesale in rebuild_slide10(), because its right-hand
    # column was a submission to-do list. Only the footer is substituted here.
    ("Repository: https://github.com/Toolstack7462/Mit208", f"Repository: {REPO}"),
]

# Slide 10, as specified for the final deck: outcome and forward look, no status.
FINAL_OUTCOME = [
    "Core analyst and staff workflow completed end to end",
    "Server-side role, ownership and state controls enforced",
    "PostgreSQL, SQLite fallback, tests and CI verified",
    "Explainable rule-based scoring and audit evidence demonstrated",
]
FUTURE_IMPROVEMENTS = [
    "Browser-level end-to-end testing",
    "Token revocation and shared rate limiting",
    "Real email-header and mailbox integration",
    "Independently evaluated classifier using a labelled dataset",
]

# Slide 4's colour-key chips, tied to the numbered layers in the new drawing.
CHIPS = {
    "Text 5": "1  BROWSER",
    "Text 7": "2  API",
    "Text 9": "3  POLICY",
    "Text 11": "4  DATA",
}

# Slide 7's metric tiles, addressed by shape name so a bare "0" cannot be replaced
# somewhere else in the deck.
TILES = {
    "Text 5": BACKEND_TESTS,
    "Text 10": COVERAGE_PCT,
    "Text 12": COVERAGE_FRAC,
    "Text 15": SMOKE,
    "Text 20": FRONTEND_TESTS,
    "Text 25": "6/6",
    "Text 26": "invalid writes",
    "Text 27": "rejected by the database",
}

NOTES: list[tuple[int, str, str]] = [
    (7, "State the figures and their environment: 118 backend tests at 89 per cent "
        "statement coverage (823 of 924 statements), run on both SQLite and PostgreSQL "
        "16.6, 69 frontend tests, and 20 live checks driving a real running server.",
        "State the figures and their environment: 170 backend tests at 90 per cent "
        "statement coverage (847 of 943 statements), run on both SQLite and PostgreSQL "
        "16.6, 92 frontend tests, and 22 live checks driving a real running server."),
    (7, "The suite defaults to in-memory SQLite and was also run in full against a "
        "PostgreSQL 16.6 server, where the schema applied cleanly, all ten check "
        "constraints and the partial unique index were confirmed and seven invalid "
        "writes were rejected.",
        "The suite defaults to in-memory SQLite and was also run in full against a "
        "PostgreSQL 16.6 server, where the schema applied cleanly, all ten check "
        "constraints and the partial unique index were confirmed and six invalid writes "
        "were rejected in raw SQL."),
    (7, "Then disclose what is still manual: confirming the public CI result and "
        "recording the narration.",
        "Then disclose what is still manual: recording the narration. The GitHub "
        "Actions run passes on the public repository, and a secret-scan job fails the "
        "build if any credential is ever committed."),
    (8, "These are captures of the running application, taken with Playwright at "
        "3200 x 2000 by evidence/capture_screenshots.py. Narrate the analyst flow and "
        "then the staff request and approval flow; do not describe every screen.",
        "These are captures of the running application, taken with Playwright at "
        "3200 x 2000 by evidence/capture_screenshots.py while it was served from "
        "PostgreSQL 16.6. Four crops rather than seven whole screens, because a whole "
        "1600px screen shrunk to a thumbnail renders its text at about 2pt. Walk the "
        "four steps in order; do not read the captions aloud."),
    (8, "Evidence: evidence/screenshots (17 images with an index).",
        "Evidence: evidence/screenshots (24 images with a generated index). The same "
        "four crops are Figure 3 in the report."),
    (10, "Close on outcome and honesty: the core MVP works end to end, is covered by "
         "179 automated tests plus 20 live checks, and is reproducible from the "
         "repository with the documented setup steps.",
         f"Close on outcome and honesty: the core MVP works end to end, is covered by "
         f"{TOTAL_TESTS} automated tests plus 22 live checks on each database engine, "
         f"and is reproducible from the repository with the documented setup steps."),
    (10, "Have the repository and the v1.0-final release open, and be ready to run the "
         "app from a terminal if the lecturer asks.",
         f"Have the repository and the {TAG} tag open, and be ready to run the app from "
         f"a terminal if the lecturer asks. docs/DEMO.md has the exact commands, "
         f"including the offline fallback."),
]

ANALYST_BULLETS = [
    "Prioritise — the inbox is ordered by risk score",
    "See why — every point has a named indicator",
    "Act safely — only the transitions the state machine allows",
    "Stay accountable — every action lands in the audit trail",
]
STAFF_BULLETS = [
    "See only mail addressed to them, enforced by the API",
    "Challenge a held message with a written justification",
    "One open request per email, per person",
    "No route that bypasses analyst review",
]
ANALYST_FILES = "routers/emails.py · transitions.py · audit.py"
STAFF_FILES = "routers/requests.py · pages/StaffPortal.jsx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def walk(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:          # group
            yield from walk(sh.shapes)


def apply_replacements(prs) -> dict[str, int]:
    hits = {old: 0 for old, _ in REPLACEMENTS}
    for slide in prs.slides:
        for sh in walk(slide.shapes):
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in REPLACEMENTS:
                        if old in run.text:
                            run.text = run.text.replace(old, new)
                            hits[old] += 1
    return hits


def apply_notes(prs) -> list[str]:
    missing = []
    for idx, old, new in NOTES:
        slide = prs.slides[idx - 1]
        if not slide.has_notes_slide:
            missing.append(f"slide {idx}: no notes slide")
            continue
        tf = slide.notes_slide.notes_text_frame
        if old not in tf.text:
            missing.append(f"slide {idx}: {old[:60]}...")
            continue
        # Notes are plain prose; rewriting the whole frame keeps paragraph breaks.
        tf.text = tf.text.replace(old, new)
    return missing


def set_shape_texts(prs, slide_index: int, mapping: dict[str, str]) -> list[str]:
    """Set the text of shapes by name, keeping their existing run formatting."""
    missing = []
    slide = prs.slides[slide_index]
    by_name = {sh.name: sh for sh in slide.shapes}
    for name, value in mapping.items():
        sh = by_name.get(name)
        if sh is None or not sh.has_text_frame or not sh.text_frame.paragraphs[0].runs:
            missing.append(name)
            continue
        runs = sh.text_frame.paragraphs[0].runs
        runs[0].text = value
        for r in runs[1:]:
            r.text = ""
    return missing


def swap_picture(slide, name: str, image: Path) -> None:
    """Replace a picture, keeping its exact frame."""
    old = next(sh for sh in slide.shapes if sh.name == name)
    left, top, width, height = old.left, old.top, old.width, old.height
    el = old._element
    el.getparent().remove(el)
    pic = slide.shapes.add_picture(str(image), left, top, width, height)
    pic.name = name
    return pic


def crop_to(src: Path, box: tuple[int, int, int, int], dst: Path) -> tuple[int, int]:
    TMP.mkdir(parents=True, exist_ok=True)
    im = Image.open(src).convert("RGB").crop(box)
    im.save(dst, "PNG")
    return im.size


def textbox(slide, left, top, width, height, text, *, size, bold=False,
            colour=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            spacing=None, font=FONT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if spacing is not None:
            para.space_after = Pt(spacing)
        run = para.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
    return box


def bulleted(para, text: str, *, size, colour=NAVY, indent_in=0.22):
    """A real PowerPoint bullet, so a wrapped line aligns under the text.

    python-pptx exposes no API for this, so the paragraph properties are set on the
    XML: marL is the text indent, a negative `indent` pulls the bullet back out to
    the margin, and buChar supplies the glyph.
    """
    pPr = para._pPr if para._pPr is not None else para._p.get_or_add_pPr()
    marL = int(Inches(indent_in))
    pPr.set("marL", str(marL))
    pPr.set("indent", str(-marL))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag.split(':')[1]}"):
            pPr.remove(el)
    bu = pPr.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar", {"char": "•"})
    pPr.append(bu)
    run = para.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = colour
    return run


def badge(slide, left, top, diameter, label: str, *, size, fill=BRAND):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top),
                                   Inches(diameter), Inches(diameter))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = label
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return shape


# ---------------------------------------------------------------------------
# Slide 2 — align the role badges and fill the empty half of each card
# ---------------------------------------------------------------------------

def rebuild_slide2(prs) -> None:
    slide = prs.slides[1]
    by_name = {sh.name: sh for sh in slide.shapes}

    # The two circles sat mid-card, vertically misaligned with each other and
    # overlapping the body text box. They become badges beside the card titles.
    for circle, letter, title, card_left in (
        ("Shape 7", "Text 8", "Text 5", 0.66),
        ("Shape 13", "Text 14", "Text 11", 8.93),
    ):
        c, l, t = by_name[circle], by_name[letter], by_name[title]
        c.left, c.top = Inches(card_left + 0.21), Inches(1.48)
        c.width = c.height = Inches(0.44)
        l.left, l.top = Inches(card_left + 0.21), Inches(1.55)
        l.width, l.height = Inches(0.44), Inches(0.30)
        for r in l.text_frame.paragraphs[0].runs:
            r.font.size = Pt(15)
        l.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        t.left, t.top = Inches(card_left + 0.78), Inches(1.47)
        t.width, t.height = Inches(2.75), Inches(0.42)

    # The body text was one sentence in a 4in-tall box, leaving the lower half of
    # each card empty. Four concrete needs each, then the files that implement them.
    for body_name, bullets, files, card_left in (
        ("Text 6", ANALYST_BULLETS, ANALYST_FILES, 0.66),
        ("Text 12", STAFF_BULLETS, STAFF_FILES, 8.93),
    ):
        body = by_name[body_name]
        body.left, body.top = Inches(card_left + 0.21), Inches(2.15)
        body.width, body.height = Inches(3.33), Inches(3.30)
        tf = body.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, line in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.space_after = Pt(11)
            bulleted(para, line, size=13)

        textbox(slide, card_left + 0.21, 5.62, 3.33, 0.50, files,
                size=9.5, colour=MUTED, spacing=0)


# ---------------------------------------------------------------------------
# Slide 8 — four readable crops instead of seven thumbnails
# ---------------------------------------------------------------------------

def rebuild_slide8(prs) -> None:
    slide = prs.slides[7]

    # Keep only the title band and the corner badge.
    keep = {"Text 0", "Shape 1", "Text 2", "Shape 3", "Text 4"}
    for sh in list(slide.shapes):
        if sh.name in keep or "Placeholder" in sh.name:
            continue
        sh._element.getparent().remove(sh._element)

    for sh in slide.shapes:
        if sh.name == "Text 2":
            sh.text_frame.paragraphs[0].runs[0].text = (
                "Screenshots from the running application, served from PostgreSQL 16.6")

    crops = []
    for i, (name, box, _, _) in enumerate(PANES, 1):
        dst = TMP / f"slide8-pane{i}.png"
        w, h = crop_to(SHOTS / name, box, dst)
        crops.append((dst, w / h))

    X0, CONTENT_W = 0.55, 12.23
    GAP_X, GAP_Y = 0.25, 0.14
    CAP_H = 0.24
    # BOTTOM leaves room for the footnote above the slide's own footer band.
    TOP, BOTTOM = 1.28, 6.52

    row1 = [c[1] for c in crops[:3]]
    h1_raw = (CONTENT_W - 2 * GAP_X) / sum(row1)
    h2_raw = CONTENT_W / crops[3][1]
    budget = (BOTTOM - TOP) - 2 * CAP_H - GAP_Y
    k = budget / (h1_raw + h2_raw)
    h1, h2 = h1_raw * k, h2_raw * k

    widths = [h1 * ar for ar in row1]
    x = X0 + (CONTENT_W - (sum(widths) + 2 * GAP_X)) / 2
    for i in range(3):
        pic = slide.shapes.add_picture(str(crops[i][0]), Inches(x), Inches(TOP),
                                       Inches(widths[i]), Inches(h1))
        pic.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        pic.line.width = Pt(0.75)
        # No badge on the image: a circle over a screenshot hides the very control
        # the pane exists to show. The caption already carries the step number.
        textbox(slide, x, TOP + h1 + 0.06, widths[i], CAP_H,
                f"{PANES[i][2]}  {PANES[i][3]}", size=10.5, bold=True, colour=MUTED)
        x += widths[i] + GAP_X

    w4 = h2 * crops[3][1]
    y4 = TOP + h1 + CAP_H + GAP_Y
    x4 = X0 + (CONTENT_W - w4) / 2
    pic = slide.shapes.add_picture(str(crops[3][0]), Inches(x4), Inches(y4),
                                   Inches(w4), Inches(h2))
    pic.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    pic.line.width = Pt(0.75)
    textbox(slide, x4, y4 + h2 + 0.06, w4, CAP_H,
            f"{PANES[3][2]}  {PANES[3][3]}", size=10.5, bold=True, colour=MUTED)

    textbox(slide, X0, 6.62, CONTENT_W, 0.26,
            "Also captured: the refused duplicate request (409), the refused short "
            "justification (422), the disabled invalid transition and the "
            "API-unreachable error state — all 24 images are in evidence/screenshots.",
            size=9.5, bold=True, colour=RGBColor(0xB6, 0x3A, 0x3A), align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 10 — conclusion and repository, with no submission status
# ---------------------------------------------------------------------------

def rebuild_slide10(prs) -> None:
    slide = prs.slides[9]

    # Keep the title band, both card backgrounds, their headings and the footer bar.
    keep = {"Text 0", "Shape 1", "Text 2", "Shape 3", "Text 4",
            "Shape 7", "Text 8", "Shape 24", "Text 25"}
    for sh in list(slide.shapes):
        if sh.name in keep or "Placeholder" in sh.name:
            continue
        sh._element.getparent().remove(sh._element)

    by_name = {sh.name: sh for sh in slide.shapes}
    heading = {
        "Text 0": "10. Conclusion and repository",
        "Text 2": "Final technical outcome, future improvements and the assessed release",
        "Text 4": "Final outcome",
        "Text 8": "Future improvements",
    }
    for name, value in heading.items():
        runs = by_name[name].text_frame.paragraphs[0].runs
        runs[0].text = value
        for r in runs[1:]:
            r.text = ""

    # The right-hand heading was sized for "Before the final submission"; the new
    # label is shorter, so leave the box and just widen nothing.
    for name, bullets, left, width in (
        ("Text 4", FINAL_OUTCOME, 1.08, 4.70),
        ("Text 8", FUTURE_IMPROVEMENTS, 6.78, 5.45),
    ):
        top = Emu(by_name[name].top).inches + 0.75
        box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(3.40))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, line in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.space_after = Pt(14)
            bulleted(para, line, size=14.5,
                     colour=WHITE if name == "Text 4" else NAVY)


# ---------------------------------------------------------------------------
# Hidden content: speaker notes, comments, metadata
# ---------------------------------------------------------------------------

def export_notes(prs, path: Path) -> int:
    """Write the speaker notes to a private file before they are stripped."""
    blocks = []
    for i, slide in enumerate(prs.slides, 1):
        if not slide.has_notes_slide:
            continue
        text = slide.notes_slide.notes_text_frame.text.strip()
        if text:
            blocks.append(f"## Slide {i}\n\n{text}\n")
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        "# PhishGuard — private rehearsal prompts\n\n"
        "Speaker notes lifted out of the deck before the submitted copy was cleaned.\n"
        "**This file is not part of the submission.** It is the presenter's own\n"
        "prompt sheet; the deck handed in carries no notes.\n\n"
        "---\n\n" + "\n---\n\n".join(blocks),
        encoding="utf-8")
    return len(blocks)


def strip_notes(prs) -> int:
    """Remove every notes slide from the presentation part graph."""
    removed = 0
    for slide in prs.slides:
        if not slide.has_notes_slide:
            continue
        part = slide.part
        for rel_id, rel in list(part.rels.items()):
            if rel.reltype.endswith("/notesSlide"):
                part.drop_rel(rel_id)
                removed += 1
    return removed


def scrub_properties(props, *, title: str) -> None:
    """Clear authorship and draft metadata on a DOCX or PPTX core-properties part."""
    props.author = ""
    props.last_modified_by = ""
    props.comments = ""
    props.category = ""
    props.keywords = ""
    props.subject = ""
    props.title = title
    props.content_status = ""
    props.identifier = ""
    props.language = "en-AU"
    props.revision = 1


# ---------------------------------------------------------------------------

def main() -> int:
    if len(sys.argv) != 3:
        raise SystemExit(__doc__)
    src, dst = Path(sys.argv[1]), Path(sys.argv[2])
    prs = Presentation(str(src))

    hits = apply_replacements(prs)
    unmatched = [old for old, n in hits.items() if n == 0]
    if unmatched:
        raise SystemExit("these strings were not found, so the deck would keep a stale "
                         "claim:\n  " + "\n  ".join(repr(u) for u in unmatched))

    missing_tiles = set_shape_texts(prs, 6, TILES)
    if missing_tiles:
        raise SystemExit(f"slide 7 tiles not found: {missing_tiles}")

    missing_chips = set_shape_texts(prs, 3, CHIPS)
    if missing_chips:
        raise SystemExit(f"slide 4 chips not found: {missing_chips}")

    missing_notes = apply_notes(prs)
    if missing_notes:
        raise SystemExit("speaker notes not found:\n  " + "\n  ".join(missing_notes))

    # Slide 1: the inbox crop was taken on 5 August; retake it from today's capture.
    TMP.mkdir(parents=True, exist_ok=True)
    hero = TMP / "slide1-hero.png"
    crop_to(SHOTS / "04-analyst-inbox.png", (560, 0, 2644, 2000), hero)
    swap_picture(prs.slides[0], "Picture 25", hero)

    # Slide 4: the simplified layered drawing, same frame and aspect ratio.
    swap_picture(prs.slides[3], "Image 0", FIGURES / "architecture-slide.png")

    rebuild_slide2(prs)
    rebuild_slide8(prs)
    rebuild_slide10(prs)

    # The notes are the presenter's own prompt sheet. They are corrected above, then
    # lifted out to a private file, then removed: a submitted deck should carry no
    # notes at all, and a marker opening the PPTX should not find rehearsal text.
    rehearsal = REHEARSAL_DIR / "SPEAKER_PROMPTS.md"
    exported = export_notes(prs, rehearsal)
    stripped = strip_notes(prs)

    scrub_properties(prs.core_properties,
                     title="PhishGuard — Technical Presentation (MIT208 Assessment 3)")

    prs.save(str(dst))
    print(f"Presentation written: {dst}")
    print(f"  {len(REPLACEMENTS)} text substitutions, all matched")
    print(f"  {len(TILES)} metric tiles updated on slide 7")
    print(f"  {len(NOTES)} speaker-note passages corrected before export")
    print("  slide 1 hero image and slide 4 architecture drawing replaced")
    print("  slides 2, 8 and 10 rebuilt")
    print(f"  {exported} slides of notes exported to {rehearsal}")
    print(f"  {stripped} notes slides removed from the submitted deck")
    print("  core properties scrubbed of authorship and draft metadata")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
