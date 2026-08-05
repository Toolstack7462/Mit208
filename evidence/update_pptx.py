"""Improve the existing PhishGuard technical presentation in place.

This does NOT recreate the deck. It keeps all ten slides, their layouts, shapes,
colours and structure, and makes four kinds of change:

  1. Replaces superseded test figures with the verified ones.
  2. Swaps the placeholder images for genuine screenshots of the running
     application, preserving each frame's exact position and size.
  3. Removes "draft / replace after final build" labelling that is no longer
     true, and rewrites the remaining-work lines so they list only what is
     genuinely still outstanding.
  4. Updates the speaker notes to match, and adds the figures a presenter needs
     to hand.

Usage:
    python evidence/update_pptx.py <input.pptx> <output.pptx>
"""
from __future__ import annotations

import copy
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Pt

HERE = Path(__file__).resolve().parent
SHOTS = HERE / "screenshots"
TMP = HERE / "_pptx_tmp"

# ---------------------------------------------------------------------------
# Verified figures (docs/TESTING.md, 5 August 2026)
# ---------------------------------------------------------------------------
BACKEND_TESTS = "110"
COVERAGE_PCT = "89%"
COVERAGE_FRAC = "805 / 909"
SMOKE = "20/20"
FRONTEND_TESTS = "69/69"
JS_FILES = "30 JS/JSX files"

# Exact text replacements, applied to every run in every shape on every slide.
TEXT_REPLACEMENTS: list[tuple[str, str]] = [
    # ---- Slide 7: test scoreboard -----------------------------------------
    ("52", BACKEND_TESTS),
    ("85%", COVERAGE_PCT),
    ("716 / 846", COVERAGE_FRAC),
    ("14/14", SMOKE),
    ("5/5", FRONTEND_TESTS),
    ("frontend policy tests", "frontend tests"),
    ("Node helpers", "vitest + jsdom"),
    ("21 JS/JSX files", JS_FILES),
    (
        "Wrong credentials, malformed/expired/wrong-type token and stored-role mismatch",
        "Wrong credentials; forged, expired and wrong-type tokens; tampered role claim",
    ),
    (
        "Duplicate quarantine, release from invalid state and stale approval",
        "Repeated no-op action, over-long input and stale approval",
    ),
    (
        "Manual final proof still required: PostgreSQL workflow + fresh npm build/browser "
        "capture + passing public CI + genuine final release.",
        "Still manual: one local PostgreSQL run (CI already tests PostgreSQL 16), "
        "confirming the public CI result, and recording the narration.",
    ),
    # ---- Slide 8: product screenshots -------------------------------------
    (
        "Current repository screenshots - replace with the final patched browser run",
        "Captured from the running application with Playwright (3200 x 2000)",
    ),
    ("REPLACE AFTER FINAL BUILD", "CAPTURED FROM THE RUNNING APP"),
    (
        "Final capture must also show a deliberate 409/validation example and the "
        "PostgreSQL database status.",
        "Also captured: refused duplicate request (409), refused short justification "
        "(422) and the API-unreachable error state.",
    ),
    # ---- Slide 9: results and limitations ---------------------------------
    (
        "Core analyst and staff journeys are implemented and API-tested",
        "Core analyst and staff journeys implemented, API-tested and screen-captured",
    ),
    (
        "Repository contains setup, schema, evidence and release plan",
        "Repository holds setup, schema, six technical documents and evidence",
    ),
    (
        "No production MFA, rate limiting, revocation or tamper-evident audit",
        "No MFA, token revocation or tamper-evident audit; rate limiting is per process",
    ),
    (
        "PostgreSQL/browser/build/CI/release evidence still manual",
        "One local PostgreSQL run and the public CI result still to confirm",
    ),
    (
        "Earlier screenshots do not prove the patched interface",
        "No automated browser-level end-to-end test yet",
    ),
    (
        "Local storage token and small-scale dashboard queries",
        "Browser-stored token, no HTTPS locally, small-scale dashboard queries",
    ),
    # ---- Slide 10: closing / next steps -----------------------------------
    ("Apply and review the patch in the real repository",
     "Confirm the passing GitHub Actions run"),
    # ---- Slide 1: title ----------------------------------------------------
    ("DRAFT - REPLACE FINAL SCREENSHOTS", "LIVE APPLICATION - CAPTURED 5 AUG 2026"),
]

# Slide 8 frames, in the order they appear, mapped to the screenshot that matches
# each numbered caption already on the slide.
SLIDE8_IMAGES = [
    "01-login.png",                            # 1 Login and role authentication
    "03-analyst-dashboard.png",                # 2 Role-aware dashboard
    "06-email-detail-explainable-score.png",   # 3 Explainable email review
    "07-audit-log.png",                        # 4 Action traceability
    "09-staff-portal.png",                     # 5 Staff held-email request
    "08-release-requests-analyst.png",         # 6 Analyst/admin decision
    "13-release-request-approved.png",         # 7 Updated status and workflow
]

# Slide 1's frame is nearly square, so the screenshot is centre-cropped to fit
# rather than stretched.
SLIDE1_IMAGE = "04-analyst-inbox.png"

# Speaker notes, rewritten to match what the deck now says. Kept to the same
# short, presenter-facing style as the originals.
NOTES: dict[int, str] = {
    1: (
        "Slide 1 - Title and project identity\n\n"
        "Introduce the project in one sentence: PhishGuard holds suspicious email "
        "instead of deleting it, shows an analyst why it was flagged, and gives the "
        "recipient a governed way to ask for it back. State that this is an individual "
        "full-stack prototype and that detection is a transparent rule engine, not a "
        "trained model.\n\n"
        "Say your name and student ID, and have the repository open in a second tab.\n\n"
        "Evidence: the image is a screenshot of the running application."
    ),
    2: (
        "Slide 2 - Problem and users\n\n"
        "Explain the two failure modes this sits between: silent deletion destroys "
        "legitimate mail with no recourse; a warning label leaves an untrained "
        "recipient making a security decision. Neither records who decided what.\n\n"
        "Cover the analyst need (prioritise, see why, act, retain a trail) and the "
        "staff need (see only their own mail, challenge a held item without bypassing "
        "review). Be explicit that samples are synthetic and there is no live mailbox."
    ),
    3: (
        "Slide 3 - Final scope and completed features\n\n"
        "Walk the completed MVP left to right, then disclose the exclusions before "
        "anyone asks. Say plainly: the DistilBERT classifier is not built - "
        "ml_model.py documents the integration point and raises NotImplementedError, "
        "and no running code path calls it. SPF, DKIM and DMARC are simulated from the "
        "rule engine's own signals because the sample data has no real SMTP headers.\n\n"
        "The end-to-end path along the bottom is exactly what the walkthrough video "
        "demonstrates."
    ),
    4: (
        "Slide 4 - Architecture\n\n"
        "Name each component: React 18 + Vite + Tailwind + React Router on the client; "
        "FastAPI with SQLAlchemy 2 and Pydantic v2 for the API; PostgreSQL as the "
        "assessed target with SQLite as a zero-install fallback.\n\n"
        "The key point: the frontend route guards are a usability feature, not a "
        "security control. Every role, ownership and state rule is re-checked in the "
        "API, and the role is re-read from the database rather than trusted from the "
        "token - so a tampered token grants nothing. Database CHECK constraints and a "
        "partial unique index are the final integrity guard.\n\n"
        "Evidence: docs/ARCHITECTURE.md and docs/ERD.md."
    ),
    5: (
        "Slide 5 - Technology stack\n\n"
        "Justify choices rather than listing them. FastAPI for typed request handling "
        "and automatic OpenAPI documentation. Pydantic to validate at the API boundary, "
        "with lengths matched to the database column widths - that one mattered, "
        "because SQLite silently accepted an over-long subject while PostgreSQL "
        "rejected it and the user saw a 500. bcrypt for hashing and PyJWT for tokens.\n\n"
        "If asked why not a trained classifier: I can explain and defend every point "
        "the rule engine awards, which matters more in a review tool than raw accuracy, "
        "and I have no labelled corpus to evaluate a model honestly."
    ),
    6: (
        "Slide 6 - Implementation journey\n\n"
        "Describe the sequence honestly: initial full-stack MVP, interface redesign, "
        "dashboard chart, then a structured review against the assessment criteria that "
        "produced the hardening work.\n\n"
        "Pick one challenge and go deep. The strongest is the dependency failure: the "
        "pinned versions could not be installed on Python 3.13 or later because "
        "pydantic-core and psycopg2 publish wheels only for the interpreters that "
        "existed at that patch release. pip fell back to a source build and failed. The "
        "fix was compatible-release pins, and a CI matrix across 3.11, 3.12 and 3.13 "
        "now stops it recurring.\n\n"
        "Evidence: docs/BUG_LOG.md, BUG-01."
    ),
    7: (
        "Slide 7 - Testing and security\n\n"
        "State the figures and their environment: 110 backend tests at 89 per cent "
        "statement coverage (805 of 909 statements), 69 frontend tests, and 20 live "
        "checks driving a real running server. The automated suite uses in-memory "
        "SQLite; a separate CI job applies the schema to a real PostgreSQL 16 service "
        "and verifies the seed.\n\n"
        "Give two representative negative tests. Good ones: a token re-signed with a "
        "tampered role claim still returns 403 because the role is re-read from the "
        "database; and a 501-character subject is refused with 422 instead of reaching "
        "the database, so SQLite and PostgreSQL behave identically.\n\n"
        "Then disclose what is still manual: one local PostgreSQL run and confirming "
        "the public CI result.\n\n"
        "Evidence: docs/TESTING.md, docs/SECURITY.md, docs/BUG_LOG.md."
    ),
    8: (
        "Slide 8 - Final product screenshots\n\n"
        "These are captures of the running application, taken with Playwright at "
        "3200 x 2000 by evidence/capture_screenshots.py. Narrate the analyst flow and "
        "then the staff request and approval flow; do not describe every screen.\n\n"
        "If you have time, point at the explainable-score panel: score 100 with six "
        "named indicators, including the display name claiming PayPal while the domain "
        "is 'paypa1' with a numeral one, and link text showing paypal.com while the "
        "href points at a raw IP address.\n\n"
        "Evidence: evidence/screenshots (17 images with an index)."
    ),
    9: (
        "Slide 9 - Results and limitations\n\n"
        "Lead with what works, then be first to name the limitations - it reads as "
        "confidence, not weakness.\n\n"
        "The one to state clearly: the rule engine is heuristic and will produce false "
        "positives and false negatives. That is precisely why the workflow quarantines "
        "rather than deletes and requires a human decision. No accuracy figure is "
        "claimed anywhere, because there is no labelled evaluation corpus.\n\n"
        "Also state: no HTTPS in the local demo, the token lives in browser storage, "
        "there is no token revocation, and the rate limiter counts within one process."
    ),
    10: (
        "Slide 10 - Conclusion and repository\n\n"
        "Close on outcome and honesty: the core MVP works end to end, is covered by "
        "179 automated tests plus 20 live checks, and is reproducible from the "
        "repository with the documented setup steps.\n\n"
        "Name the realistic next steps: a browser-level end-to-end suite, token "
        "revocation and a shared rate-limit store, database enumerated types, and only "
        "then a trained classifier with a documented dataset and error analysis.\n\n"
        "Have the repository and the v1.0-final release open, and be ready to run the "
        "app from a terminal if the lecturer asks."
    ),
}


def replace_text_everywhere(prs, pairs: list[tuple[str, str]]) -> int:
    """Replace text run by run so fonts, sizes and colours are preserved."""
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for old, new in pairs:
                        # Exact-match short numeric strings so "52" cannot corrupt
                        # a longer number elsewhere in the deck.
                        if len(old) <= 5:
                            if run.text.strip() == old:
                                run.text = run.text.replace(old, new)
                                count += 1
                        elif old in run.text:
                            run.text = run.text.replace(old, new)
                            count += 1
    return count


def fit_image(src: Path, target_ratio: float, out: Path) -> Path:
    """Centre-crop src to target_ratio (width/height) without stretching."""
    with Image.open(src) as im:
        w, h = im.size
        current = w / h
        if abs(current - target_ratio) < 0.01:
            return src
        if current > target_ratio:            # too wide -> crop sides
            new_w = int(h * target_ratio)
            left = (w - new_w) // 2
            box = (left, 0, left + new_w, h)
        else:                                  # too tall -> crop top/bottom
            new_h = int(w / target_ratio)
            top = (h - new_h) // 2
            box = (0, top, w, top + new_h)
        im.crop(box).save(out)
    return out


def swap_picture(slide, picture, image_path: Path) -> None:
    """Replace a picture's image while keeping its frame exactly where it is.

    python-pptx has no public API for this, so the new picture is inserted at the
    old one's position and size, then moved to the old one's z-order position and
    the old one removed. That preserves any shape drawn on top of the frame.
    """
    left, top, width, height = picture.left, picture.top, picture.width, picture.height
    old_el = picture._element
    parent = old_el.getparent()
    index = list(parent).index(old_el)

    new_pic = slide.shapes.add_picture(str(image_path), left, top, width, height)
    new_el = new_pic._element
    parent.remove(new_el)
    parent.insert(index, new_el)
    parent.remove(old_el)


def set_notes(slide, text: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    lines = text.split("\n")
    tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        tf.add_paragraph().text = line
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(12)


def main() -> int:
    if len(sys.argv) != 3:
        print(__doc__)
        return 1
    src, dst = Path(sys.argv[1]), Path(sys.argv[2])
    if not SHOTS.is_dir():
        print(f"Screenshots not found at {SHOTS}. Run capture_screenshots.py first.")
        return 1
    TMP.mkdir(exist_ok=True)

    prs = Presentation(str(src))

    # 1. Figures and wording -------------------------------------------------
    changed = replace_text_everywhere(prs, TEXT_REPLACEMENTS)
    print(f"Text runs updated: {changed}")

    # 2. Slide 1 hero image --------------------------------------------------
    slide1 = prs.slides[0]
    pics1 = [s for s in slide1.shapes if s.shape_type == 13]
    if pics1:
        pic = pics1[0]
        ratio = pic.width / pic.height
        fitted = fit_image(SHOTS / SLIDE1_IMAGE, ratio, TMP / "slide1.png")
        swap_picture(slide1, pic, fitted)
        print(f"Slide 1: image replaced with {SLIDE1_IMAGE} "
              f"(centre-cropped to {ratio:.2f}:1)")

    # 3. Slide 8 screenshot grid --------------------------------------------
    slide8 = prs.slides[7]
    pics8 = [s for s in slide8.shapes if s.shape_type == 13]
    print(f"Slide 8: {len(pics8)} image frames, {len(SLIDE8_IMAGES)} screenshots to place")
    for i, pic in enumerate(pics8):
        if i >= len(SLIDE8_IMAGES):
            break
        name = SLIDE8_IMAGES[i]
        source = SHOTS / name
        if not source.exists():
            print(f"  MISSING {name} - frame {i + 1} left unchanged")
            continue
        ratio = pic.width / pic.height
        fitted = fit_image(source, ratio, TMP / f"slide8_{i}.png")
        swap_picture(slide8, pic, fitted)
        print(f"  frame {i + 1}: {name}")

    # 4. Speaker notes -------------------------------------------------------
    for number, text in NOTES.items():
        set_notes(prs.slides[number - 1], text)
    print(f"Speaker notes rewritten on {len(NOTES)} slides")

    prs.save(str(dst))

    # Report any remaining draft language.
    check = Presentation(str(dst))
    flags = ["DRAFT", "REPLACE", "placeholder", "patch", "candidate", "52 backend"]
    found = []
    for i, slide in enumerate(check.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for word in flags:
                    if word.lower() in shape.text_frame.text.lower():
                        found.append((i, word, shape.text_frame.text[:70]))
    print(f"\nSaved: {dst}")
    print(f"Slides: {len(check.slides)}  |  "
          f"notes on {sum(1 for s in check.slides if s.has_notes_slide)} slides")
    if found:
        print("\nRemaining draft-language mentions to review:")
        for i, word, text in found:
            print(f"  slide {i}: '{word}' in {text!r}")
    else:
        print("No draft/placeholder language remains in slide text.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
